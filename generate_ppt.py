import sys
import subprocess

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Set presentation slide width & height (16:9 widescreen)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define color palette (inspired by HireAI theme)
DARK_BG = RGBColor(15, 23, 42)       # Slate 900
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_LIGHT = RGBColor(255, 255, 255)
PRIMARY = RGBColor(99, 102, 241)     # Indigo 500
MUTED = RGBColor(100, 116, 139)       # Slate 500

# Helper to add standard slide layout
def create_slide(title_text):
    slide_layout = prs.slide_layouts[6] # blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color to Light Slate
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Top Bar Header (inspired by redrob/H2S branding)
    # Background shape for header
    shape = slide.shapes.add_shape(
        1, # rectangle
        Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    
    # Header text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "Segoe UI"
    
    # Add branding tag on right of top bar
    txBox_brand = slide.shapes.add_textbox(Inches(10.0), Inches(0.2), Inches(2.8), Inches(0.6))
    tf_brand = txBox_brand.text_frame
    p_brand = tf_brand.paragraphs[0]
    p_brand.text = "redrob | H2S | INDIA.RUNS"
    p_brand.font.size = Pt(14)
    p_brand.font.italic = True
    p_brand.font.color.rgb = PRIMARY
    p_brand.font.name = "Segoe UI"
    p_brand.alignment = 2 # Right align
    
    # Bottom accent line (Purple)
    accent = slide.shapes.add_shape(
        1,
        Inches(0), Inches(7.3), Inches(13.333), Inches(0.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()
    
    return slide

# Slide 1: Title Slide (Dark Theme)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK_BG

# Accent shape on title slide
accent_shape = slide.shapes.add_shape(
    1,
    Inches(0), Inches(7.1), Inches(13.333), Inches(0.4)
)
accent_shape.fill.solid()
accent_shape.fill.fore_color.rgb = PRIMARY
accent_shape.line.fill.background()

# Title text
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "HireAI"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.font.name = "Segoe UI"

p2 = tf.add_paragraph()
p2.text = "AI-Powered Candidate Ranking & Semantic Alignment Platform"
p2.font.size = Pt(24)
p2.font.color.rgb = TEXT_LIGHT
p2.font.name = "Segoe UI"

# Info box
txBox_info = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(8.0), Inches(2.5))
tf_info = txBox_info.text_frame
def add_info_line(label, value):
    p = tf_info.add_paragraph()
    run1 = p.add_run()
    run1.text = label + ": "
    run1.font.bold = True
    run1.font.color.rgb = MUTED
    run1.font.size = Pt(16)
    run2 = p.add_run()
    run2.text = value
    run2.font.color.rgb = TEXT_LIGHT
    run2.font.size = Pt(16)

add_info_line("Project Name", "HireAI")
add_info_line("Hackathon Track", "INDIA.RUNS (redrob | H2S)")
add_info_line("Developer", "Solo Participant")

# Slide 2: Solution Overview
slide2 = create_slide("Solution Overview")
txBox = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "What is your proposed solution?"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = PRIMARY

p_ans = tf.add_paragraph()
p_ans.text = "HireAI is a next-generation candidate matching and evaluation copilot for recruiters. It parses resumes and job descriptions, runs a weighted multi-signal ranking engine, identifies and visualizes skill gaps, and includes an interactive AI Copilot chat to dynamically query the candidate database."
p_ans.font.size = Pt(16)
p_ans.font.color.rgb = TEXT_DARK
p_ans.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "What differentiates your approach from traditional candidate matching systems?"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = PRIMARY

p2_ans = tf.add_paragraph()
p2_ans.text = "• Semantic Search vs. Keywords: Traditional systems rely on exact string matches. HireAI uses vector embeddings to match concepts and synonyms.\n• Interactive Explainability: Rather than a black-box score, recruiters can talk to the Copilot to get side-by-side candidate comparisons and explanations.\n• Visual Skill-Gap Analysis: Auto-highlights missing and matched skills with integrated visualizations."
p2_ans.font.size = Pt(16)
p2_ans.font.color.rgb = TEXT_DARK

# Slide 3: JD Understanding & Candidate Evaluation
slide3 = create_slide("JD Understanding & Candidate Evaluation")
tx = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "What are the key requirements extracted from the JD?"
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = PRIMARY

p1_ans = tf.add_paragraph()
p1_ans.text = "• Required & Preferred Skills: Automatically parsed to construct target lists.\n• Experience Thresholds: Minimum and maximum years of experience extracted to calculate constraints.\n• Contextual Raw Text: Sent to the embedding engine to capture implicit responsibilities."
p1_ans.font.size = Pt(16)
p1_ans.font.color.rgb = TEXT_DARK
p1_ans.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "Which candidate signals are most important? How does the solution go beyond keywords?"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = PRIMARY

p2_ans = tf.add_paragraph()
p2_ans.text = "• Cosine Similarity (45%): Dense vectors encode candidates' summaries, experience, and projects. Cosine similarity calculates deep contextual fit.\n• Fuzzy Skill Alignment (25%): Uses token-sort ratio algorithms (RapidFuzz) to align candidate skills with JD requirements, matching spelling variants (e.g., 'ReactJS' vs 'React')."
p2_ans.font.size = Pt(16)
p2_ans.font.color.rgb = TEXT_DARK

# Slide 4: Ranking Methodology
slide4 = create_slide("Ranking Methodology")
tx = slide4.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "How does the system retrieve, score, and rank candidates?"
p1.font.size = Pt(18)
p1.font.bold = True
p1.font.color.rgb = PRIMARY

p1_ans = tf.add_paragraph()
p1_ans.text = "Resumes are parsed via pdfplumber. Profiles are compiled into structured SQLite tables. When a JD is activated, the ranking engine runs locally, scoring every candidate against the JD across four pillars, storing the results in database tables to immediately serve a React frontend leaderboard."
p1_ans.font.size = Pt(14)
p1_ans.font.color.rgb = TEXT_DARK
p1_ans.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "What models, algorithms, or heuristics are used?"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = PRIMARY

p2_ans = tf.add_paragraph()
p2_ans.text = "• Model: Hugging Face 'all-MiniLM-L6-v2' SentenceTransformer (384-dimensional embeddings).\n• Fuzzy Algorithm: RapidFuzz Levenshtein string matching for skills.\n• Heuristics: Boundary penalty calculations for experience ranges & score mappings for behavioral parameters."
p2_ans.font.size = Pt(14)
p2_ans.font.color.rgb = TEXT_DARK
p2_ans.space_after = Pt(10)

p3 = tf.add_paragraph()
p3.text = "How are multiple candidate signals combined into a final ranking?"
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = PRIMARY

p3_ans = tf.add_paragraph()
p3_ans.text = "Through a composite weighted scoring formula:\nFinal Score = (Semantic Similarity * 45%) + (Fuzzy Skill Match * 25%) + (Behavioral Score * 20%) + (Experience Match * 10%)"
p3_ans.font.size = Pt(14)
p3_ans.font.color.rgb = TEXT_DARK

# Slide 5: Explainability & Data Validation
slide5 = create_slide("Explainability & Data Validation")
tx = slide5.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "How are ranking decisions explained? How do you prevent hallucinations?"
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = PRIMARY

p1_ans = tf.add_paragraph()
p1_ans.text = "• Ground-Truth Explanations: Explainability does not rely on LLM prompts that can hallucinate. Instead, it queries the SQLite DB for actual computed sub-scores, matched skills, and missing skills to display precise bulleted strengths.\n• Chat Verification: Copilot reports exact database values in markdown comparisons (e.g. 'Priya has 95% semantic match while Sofia has 55%')."
p1_ans.font.size = Pt(16)
p1_ans.font.color.rgb = TEXT_DARK
p1_ans.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "How does your solution handle inconsistent or low-quality profiles?"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = PRIMARY

p2_ans = tf.add_paragraph()
p2_ans.text = "• Profile Completeness: Evaluates completeness based on missing critical fields (summary, history, education, projects).\n• Safe Fallbacks: Gracefully handles empty fields and JSON arrays, enabling ranking computations to run even with partially corrupted or low-text resume uploads."
p2_ans.font.size = Pt(16)
p2_ans.font.color.rgb = TEXT_DARK

# Slide 6: End-to-End Workflow
slide6 = create_slide("End-to-End Workflow")
tx = slide6.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Workflow Architecture: From JD Input to Leaderboard Export"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.space_after = Pt(15)

p_flow = tf.add_paragraph()
p_flow.text = "1. Upload JD: The recruiter pastes a job description. The FastAPI backend parses the required skills and experience thresholds.\n\n2. Parse Resumes: The recruiter uploads PDF/DOCX resumes. The text is parsed via pdfplumber and structured into SQLite database entries.\n\n3. Compute Vector Embeddings: The profile texts are converted into 384-dimensional dense vectors via the SentenceTransformer engine.\n\n4. Execute Ranking: The Weighted Composite Engine runs in sub-seconds to compute Semantic, Skill, Experience, and Behavioral metrics.\n\n5. Review & Export: The recruiter views the ranked leaderboard, reviews missing skills on the visual dashboard, chats with the Copilot, and exports candidates as a CSV file."
p_flow.font.size = Pt(16)
p_flow.font.color.rgb = TEXT_DARK

# Slide 7: System Architecture
slide7 = create_slide("System Architecture")
tx = slide7.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Technology Stack & Layered Design"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.space_after = Pt(15)

p_arch = tf.add_paragraph()
p_arch.text = "• Frontend Presentation Layer: Developed in React.js and TypeScript, optimized with Vite for hot-reloads. Recharts handles interactive visualizations (KPI dashboard, candidate match graphs). Styling is built with Tailwind CSS v4.\n\n• Backend Logic Layer: Built with FastAPI (Python) for asynchronous endpoints. Integrates SQLAlchemy ORM to manage SQLite data schemas.\n\n• Machine Learning & Match Layer: Local inference using Hugging Face Sentence Transformers (all-MiniLM-L6-v2) for semantic vector alignments. RapidFuzz handles string distance ratio calculations.\n\n• Hosting & Infrastructure: Deployed live via Render web services, featuring a containerized Docker build that pre-caches models to ensure sub-second startup."
p_arch.font.size = Pt(16)
p_arch.font.color.rgb = TEXT_DARK

# Slide 8: Results & Performance
slide8 = create_slide("Results & Performance")
tx = slide8.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "What results or insights demonstrate ranking quality?"
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = PRIMARY

p1_ans = tf.add_paragraph()
p1_ans.text = "• Synonyms Matching: Demonstrates that candidate alignment scores correctly identify conceptual fits (e.g. matching NLP experience to LLMs) where keyword-based ATS fail.\n• Holistic Ranking: The composite weighted metric ensures candidates with strong profiles and high response rates outrank inactive or slightly over-qualified candidates, reflecting real-world recruitment criteria."
p1_ans.font.size = Pt(16)
p1_ans.font.color.rgb = TEXT_DARK
p1_ans.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "How does your solution meet the challenge's runtime and compute constraints?"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = PRIMARY

p2_ans = tf.add_paragraph()
p2_ans.text = "• Pre-caching Models: Docker image pre-caches MiniLM-L6-v2 embeddings, removing download latency during cold starts.\n• Sub-second Execution: Local inference and fuzzy calculations are completed in <500ms for dozens of profiles, meeting strict real-time constraint limits."
p2_ans.font.size = Pt(16)
p2_ans.font.color.rgb = TEXT_DARK

# Save presentation
prs.save("presentation.pptx")
print("Presentation generated successfully as presentation.pptx!")
